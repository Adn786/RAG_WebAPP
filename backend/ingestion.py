"""File parsing (per extension) and token-based chunking."""
import docx
import pandas as pd
import PyPDF2
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation

from utils import get_extension, new_id, setup_logger, sha256_hash, string_id_to_int64

logger = setup_logger(__name__)


def _df_to_markdown(df: pd.DataFrame) -> str:
    try:
        return df.to_markdown(index=False)
    except Exception:
        return df.to_csv(index=False)


def parse_pdf(file_path):
    sections = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                sections.append({"text": text, "page": i, "section": None})
    return sections


def parse_docx(file_path):
    sections = []
    document = docx.Document(file_path)

    body_text = [p.text for p in document.paragraphs if p.text.strip()]
    if body_text:
        sections.append({"text": "\n".join(body_text), "page": None, "section": "body"})

    for t_idx, table in enumerate(document.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows) > 1 else pd.DataFrame(rows)
        sections.append({"text": _df_to_markdown(df), "page": None, "section": f"table_{t_idx}"})
    return sections


def parse_pptx(file_path):
    sections = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                if rows:
                    df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows) > 1 else pd.DataFrame(rows)
                    texts.append(_df_to_markdown(df))
        if texts:
            sections.append({"text": "\n".join(texts), "page": i, "section": f"slide_{i}"})
    return sections


def parse_excel_csv(file_path, extension):
    sections = []
    if extension == "csv":
        df = pd.read_csv(file_path)
        sections.append({"text": _df_to_markdown(df), "page": None, "section": "sheet_1"})
    else:
        sheets = pd.read_excel(file_path, sheet_name=None)
        for name, df in sheets.items():
            sections.append({"text": _df_to_markdown(df), "page": None, "section": name})
    return sections


def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [{"text": text, "page": None, "section": None}] if text.strip() else []


def parse_file(file_path, filename):
    """Dispatches to the right parser based on file extension. Returns a list
    of section dicts: {text, page, section}. Tables are rendered as markdown."""
    extension = get_extension(filename)
    if extension == "pdf":
        sections = parse_pdf(file_path)
    elif extension == "docx":
        sections = parse_docx(file_path)
    elif extension == "pptx":
        sections = parse_pptx(file_path)
    elif extension in ("xlsx", "csv"):
        sections = parse_excel_csv(file_path, extension)
    elif extension == "txt":
        sections = parse_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: .{extension}")

    if not sections:
        raise ValueError("No extractable text found in document")
    return sections


def chunk_sections(sections, doc_id, source, chunk_size=512, chunk_overlap=50):
    """Splits each section's text into token-based chunks (RecursiveCharacterTextSplitter
    over the cl100k_base tiktoken encoding), preserving page/section metadata."""
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )

    chunks = []
    for section in sections:
        for piece in splitter.split_text(section["text"]):
            if not piece.strip():
                continue
            chunk_id = new_id()
            chunks.append({
                "id": chunk_id,
                "int_id": string_id_to_int64(chunk_id),
                "doc_id": doc_id,
                "text": piece,
                "page": section.get("page"),
                "source": source,
                "section": section.get("section"),
                "embedding_hash": sha256_hash(piece),
            })
    return chunks
